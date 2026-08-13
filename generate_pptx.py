#!/usr/bin/env python3
"""
Generate a 12-slide PPTX presentation for the ABC Retail case study.
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE   = Path(__file__).parent
CHARTS = BASE / 'output' / 'charts'
OUT    = BASE / 'output'

# Load model results
with open(OUT / 'results_summary.json') as f:
    R = json.load(f)

# Colours
NAVY    = RGBColor(0x0B, 0x1F, 0x3A)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0x94, 0xA3, 0xB8)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
RED     = RGBColor(0xEF, 0x44, 0x44)
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
BG      = RGBColor(0xFA, 0xFB, 0xFD)
DIVIDER = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

def add_title_bar(slide, text, subtitle=None):
    # Top navy bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY; s.line.fill.background()
    tf = s.text_frame; tf.clear(); tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(0.15)
    run = tf.paragraphs[0].add_run()
    run.text = text; run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = subtitle; r2.font.size = Pt(14); r2.font.color.rgb = LGRAY

def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return tf

def add_bullet_list(slide, left, top, width, height, items, size=13, color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0; p.space_after = Pt(6)
        run = p.add_run()
        run.text = item; run.font.size = Pt(size); run.font.color.rgb = color
    return tf

def add_kpi_card(slide, left, top, label, value, sub, accent_color):
    w, h = Inches(2.6), Inches(1.5)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = DIVIDER; s.line.width = Pt(1)
    # Accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent_color; bar.line.fill.background()
    # Label
    add_text_box(slide, left+Inches(0.15), top+Inches(0.15), w-Inches(0.3), Inches(0.3), label, size=9, color=LGRAY, bold=True)
    # Value
    add_text_box(slide, left+Inches(0.15), top+Inches(0.5), w-Inches(0.3), Inches(0.5), value, size=22, color=accent_color, bold=True)
    # Sub
    add_text_box(slide, left+Inches(0.15), top+Inches(1.05), w-Inches(0.3), Inches(0.4), sub, size=8, color=LGRAY)

def add_image(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)

# ------------------------------------------------------------------
# Slide 1: Title
# ------------------------------------------------------------------
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
# Full navy background for title
bg_shape = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = NAVY; bg_shape.line.fill.background()
add_text_box(s1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
             'ABC Retail', size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s1, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             'Marketing Channel Effectiveness & Budget Optimisation', size=24, color=LGRAY, align=PP_ALIGN.CENTER)
# Divider
div = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = BLUE; div.line.fill.background()
add_text_box(s1, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
             'Marketing Mix Modeling  |  Data-Driven Budget Allocation', size=14, color=LGRAY, align=PP_ALIGN.CENTER)
add_text_box(s1, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             'Abinash Sahu  |  Team Lead — Data Intelligence  |  August 2026', size=13, color=LGRAY, align=PP_ALIGN.CENTER)
add_text_box(s1, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             'Marketing Mix Modeling Case Study', size=12, color=LGRAY, align=PP_ALIGN.CENTER)

# ------------------------------------------------------------------
# Slide 2: Executive Summary
# ------------------------------------------------------------------
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2); add_title_bar(s2, 'Executive Summary', '36 months of sales and marketing data, Jan 2020 – Dec 2022')

add_kpi_card(s2, Inches(0.7),  Inches(1.6), 'TOTAL REVENUE', '$8,750M', '3-year aggregate', BLUE)
add_kpi_card(s2, Inches(3.5),  Inches(1.6), 'MODEL R-SQUARED', f'{R["model"]["r2"]:.1%}', f'AIC {R["model"]["aic"]}', GREEN)
add_kpi_card(s2, Inches(6.3),  Inches(1.6), 'TOP CHANNEL', 'Digital', 'Strongest positive driver', TEAL)
add_kpi_card(s2, Inches(9.1),  Inches(1.6), 'OPTIMISATION UPLIFT', f'+{R["optimisation"]["uplift_pct"]}%', 'From budget reallocation', ORANGE)

findings = [
    'Q1: Only Digital and Radio are effective. TV has no measurable impact on sales.',
    'Q2: Digital is the most important channel (std coef 0.85, p < 0.001).',
    'Q3: GDP, CPI, and unemployment are not significant after marketing controls.',
    'Q4: COVID did not structurally damage sales — marketing spend sustained demand.',
    f'Recommendation: Reallocate budget for +{R["optimisation"]["uplift_pct"]}% media-effect uplift.',
]
add_bullet_list(s2, Inches(0.7), Inches(3.6), Inches(11.5), Inches(3.5), findings, size=15)

# ------------------------------------------------------------------
# Slide 3: Data & Methodology
# ------------------------------------------------------------------
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3); add_title_bar(s3, 'Methodology', 'Marketing Mix Modeling — industry-standard for channel evaluation')

method_items = [
    'Adstock transformation: geometric decay to capture carryover effects (e.g., a TV ad impression persists for weeks).',
    f'Decay rates found via grid search (900 combinations, minimised AIC): Digital λ={R["adstock_decay_rates"]["digital"]}, Radio λ={R["adstock_decay_rates"]["radio"]}, TV λ={R["adstock_decay_rates"]["tv"]}.',
    'Diminishing returns: log(1+x) saturation models decreasing marginal effectiveness at higher spend levels.',
    'Robust inference: HC3 standard errors handle heteroscedasticity; standard errors remain valid with small samples.',
    'Cross-validation: Leave-One-Out CV with Ridge regression — tests predictive accuracy on unseen data.',
    'Additional macro signals: CPI (inflation) and unemployment rate alongside GDP, sourced from BLS.gov.',
    'Confounders: seasonality (sin/cos), time trend, COVID lockdown and relaxation phase dummies.',
]
add_bullet_list(s3, Inches(0.7), Inches(1.7), Inches(5.8), Inches(5.5), method_items, size=12)
add_image(s3, CHARTS/'06_adstock_decay.png', Inches(7.0), Inches(1.7), width=Inches(5.8))

# ------------------------------------------------------------------
# Slide 4: Channel Effectiveness — Coefficients
# ------------------------------------------------------------------
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4); add_title_bar(s4, 'Channel Effectiveness', 'Standardised coefficients — who drives sales?')
add_image(s4, CHARTS/'07_standardized_coefficients.png', Inches(0.5), Inches(1.6), width=Inches(6.5))

ch_items = [
    'Digital: standardised β = 0.85, p < 0.001 — the strongest sales driver by a significant margin.',
    'Radio: standardised β = 0.56, p < 0.001 — second strongest, and potentially under-invested.',
    'TV: β = –0.27, p = 0.19 — not significant. Possibly too small a budget to detect, or genuinely ineffective at this scale.',
    'None of the macro or COVID variables are significant after controlling for media spend.',
]
add_bullet_list(s4, Inches(7.3), Inches(1.8), Inches(5.5), Inches(5), ch_items, size=13)

# ------------------------------------------------------------------
# Slide 5: Digital Deep-Dive
# ------------------------------------------------------------------
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5); add_title_bar(s5, 'Digital — The Growth Engine', 'Highest impact, proven ROI, but approaching saturation')
add_image(s5, CHARTS/'16_digital_saturation.png', Inches(0.5), Inches(1.6), width=Inches(6.2))
add_image(s5, CHARTS/'05_channel_scatter.png', Inches(6.8), Inches(1.6), width=Inches(6.0))

dig_items = [
    f'Marginal ROI: ${R["channels"][0]["marginal_roi"]:.2f} additional sales per $1 of Digital spend (at current levels).',
    'Strong correlation with sales (r = 0.77) — consistent month after month.',
    'Saturation curve shows diminishing returns above ~$100M/month — indicates room to redistribute excess.',
    'Fast decay (λ = 0.15) — Digital is a "pull" channel with immediate response. Requires consistent investment.',
]
add_bullet_list(s5, Inches(0.7), Inches(5.3), Inches(12), Inches(2), dig_items, size=12)

# ------------------------------------------------------------------
# Slide 6: Radio & TV
# ------------------------------------------------------------------
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6); add_title_bar(s6, 'Radio & TV Assessment', 'Radio is under-invested; TV has no detectable effect')
add_image(s6, CHARTS/'08_channel_roi.png', Inches(0.5), Inches(1.6), width=Inches(8))

rt_items = [
    f'Radio marginal ROI: ${R["channels"][1]["marginal_roi"]:.2f}/$1 — higher than Digital because it is currently under-invested (only 13% of budget).',
    'The optimizer recommends increasing Radio share from 13% to ~42% — subject to incrementality testing.',
    'TV spend is measured in thousands (not millions) — effectively a rounding error in the budget (0.02%).',
    'At this spend level, TV cannot be evaluated reliably. Recommendation: hold and measure before scaling.',
]
add_bullet_list(s6, Inches(0.5), Inches(5.2), Inches(12), Inches(2), rt_items, size=12)

# ------------------------------------------------------------------
# Slide 7: COVID Impact
# ------------------------------------------------------------------
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7); add_title_bar(s7, 'COVID-19 Impact Analysis', 'Lockdown did not collapse sales — marketing sustained demand')
add_image(s7, CHARTS/'09_covid_impact_bridge.png', Inches(0.3), Inches(1.5), width=Inches(7.5))

covid_items = [
    f'Lockdown coefficient: +$11.5M, p = {R["macro"]["covid_lockdown_p"]:.3f} — not significant.',
    f'Relaxation coefficient: –$80.1M, p = {R["macro"]["covid_relaxation_p"]:.3f} — not significant.',
    'ABC increased digital spend during lockdown ($57M vs $49M pre-COVID), offsetting reduced foot traffic.',
    'The real sales dip (Q3 2020, $179M avg) coincided with marketing budget cuts, not COVID restrictions.',
    'Implication: ABC controls its own sales trajectory through marketing. COVID was not a structural break.',
]
add_bullet_list(s7, Inches(8.0), Inches(1.8), Inches(4.8), Inches(5.2), covid_items, size=12)

# ------------------------------------------------------------------
# Slide 8: Macro-Economic Analysis
# ------------------------------------------------------------------
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8); add_title_bar(s8, 'Macro-Economic Factors', 'GDP, CPI, CA Unemployment, Sentiment — none significant after media controls')
add_image(s8, CHARTS/'02_quarterly_sales_gdp.png', Inches(0.3), Inches(1.5), width=Inches(7.5))

macro_items = [
    f'GDP quarterly change: p = {R["macro"]["gdp_p"]:.3f}. No significant effect despite Q2 2020 drop (–29.9% GDP).',
    f'CPI YoY (Inflation): p = {R["macro"]["cpi_p"]:.3f}. Inflation peaked at 9.1% — no sales impact detected.',
    f'CA Unemployment: p = {R["macro"]["ca_unemployment_p"]:.3f}. CA peak of 16.4% in May 2020 did not alter demand.',
    f'Consumer Sentiment & E-Commerce %: Sentiment (p={R["macro"]["consumer_sentiment_p"]:.3f}) & E-Commerce % (p={R["macro"]["ecommerce_pct_p"]:.3f}) are also not significant after media controls.',
    'Bay Area market insight: ABC sales are driven by marketing execution and brand awareness in target demographic (12-45 yrs), not macro cycles.',
]
add_bullet_list(s8, Inches(8.0), Inches(1.8), Inches(4.8), Inches(5.2), macro_items, size=11)

# ------------------------------------------------------------------
# Slide 9: Budget Optimisation
# ------------------------------------------------------------------
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9); add_title_bar(s9, 'Budget Optimisation', f'Constrained reallocation — predicted uplift {R["optimisation"]["uplift_pct"]:+.1f}%')
add_image(s9, CHARTS/'10_budget_allocation.png', Inches(0.3), Inches(1.5), width=Inches(7.5))

opt_items = [
    f'Current: Digital {R["optimisation"]["current"]["Digital"]/30.33:.0f}%, Radio {R["optimisation"]["current"]["Radio"]/30.33:.0f}%, TV ~0%.',
    f'Optimal: Digital {R["optimisation"]["optimal"]["Digital"]/30.33:.0f}%, Radio {R["optimisation"]["optimal"]["Radio"]/30.33:.0f}%, TV 0%.',
    'The shift reflects diminishing returns on Digital — at current high spend, each incremental dollar yields less.',
    'Radio is under-saturated: high marginal ROI ($7.07/$1) means reallocating here creates more value per dollar.',
    'This model recommendation delivers +31.7% uplift from the same $3,033M budget. Validate via geo holdout tests.',
]
add_bullet_list(s9, Inches(8.0), Inches(1.8), Inches(4.8), Inches(5.2), opt_items, size=11)

# ------------------------------------------------------------------
# Slide 10: Sensitivity
# ------------------------------------------------------------------
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s10); add_title_bar(s10, 'Sensitivity Analysis', 'Budget elasticity and scenario planning')
add_image(s10, CHARTS/'11_sensitivity_tornado.png', Inches(0.3), Inches(1.5), width=Inches(7))

sens_items = [
    'Increasing total budget by 30% yields diminishing incremental returns — consistent with saturation.',
    'Cutting budget by 20% results in meaningful loss, confirming marketing is a revenue driver.',
    'The optimal mix (58/42 Digital/Radio) is stable across budget scenarios.',
    'This sensitivity analysis helps leadership evaluate trade-offs in budget negotiation.',
]
add_bullet_list(s10, Inches(7.5), Inches(1.8), Inches(5.3), Inches(5.2), sens_items, size=13)

# ------------------------------------------------------------------
# Slide 11: Roadmap
# ------------------------------------------------------------------
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s11); add_title_bar(s11, 'Leadership Action Roadmap', '30 / 60 / 90 day plan')

phases = [
    ('30 DAYS — Foundation', [
        'Validate channel data definitions and units.',
        'Establish repeatable monthly data pipeline.',
        'Set up Digital A/B testing framework.',
        'Create KPI dictionary and reporting cadence.',
    ]),
    ('60 DAYS — Validation', [
        'Run Digital and Radio incrementality tests (geo holdout or matched-market).',
        'Add pricing, promotions, holidays, and competitor signals to the model.',
        'Begin shift of 10–15% budget from Digital to Radio as a pilot.',
    ]),
    ('90 DAYS — Scale', [
        'Deploy production MMM with automated monthly refresh.',
        'Build budget allocation engine with approval workflow.',
        'Establish model governance: drift monitoring, version control.',
        'Present quarterly business review with updated ROI estimates.',
    ]),
]

y_pos = Inches(1.7)
for title, items in phases:
    add_text_box(s11, Inches(0.8), y_pos, Inches(11), Inches(0.4), title, size=16, bold=True, color=BLUE)
    y_pos += Inches(0.4)
    add_bullet_list(s11, Inches(1.2), y_pos, Inches(10.5), Inches(len(items)*0.4), items, size=12, color=NAVY)
    y_pos += Inches(len(items)*0.35 + 0.15)

# ------------------------------------------------------------------
# Slide 12: Appendix
# ------------------------------------------------------------------
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s12); add_title_bar(s12, 'Appendix — Model Diagnostics', 'Residuals, VIF, and cross-validation results')
add_image(s12, CHARTS/'13_residual_diagnostics.png', Inches(0.3), Inches(1.5), width=Inches(7.5))

app_items = [
    f'R²: {R["model"]["r2"]:.3f}  |  Adj R²: {R["model"]["adj_r2"]:.3f}  |  LOOCV R²: {R["model"]["loocv_r2"]:.3f}',
    f'AIC: {R["model"]["aic"]}  |  BIC: {R["model"]["bic"]}',
    f'Durbin-Watson: {R["diagnostics"]["durbin_watson"]:.3f}',
    f'Breusch-Pagan p-value: {R["diagnostics"]["bp_pval"]:.4f} (homoscedastic)',
    'HC3 robust standard errors used throughout — valid under heteroscedasticity.',
    'Leave-One-Out CV with Ridge (alpha=10) — unbiased out-of-sample estimate.',
    f'Grid search tested {R["grid_search"]["combinations_tested"]} adstock decay combinations.',
    'Caveats: n=36, no geography/pricing/competitor data. Treat as test-and-learn priorities.',
]
add_bullet_list(s12, Inches(8.0), Inches(1.8), Inches(4.8), Inches(5.2), app_items, size=11)

# ------------------------------------------------------------------
# Save
# ------------------------------------------------------------------
pptx_path = BASE / 'ABC_Retail_Case_Study_Presentation.pptx'
prs.save(str(pptx_path))
print(f"Presentation saved: {pptx_path}")
print(f"  {len(prs.slides)} slides")
